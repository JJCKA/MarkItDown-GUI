"""File conversion wrapper around markitdown."""

import asyncio
import re
import tempfile
import time
from pathlib import Path
from typing import Callable

from markitdown import MarkItDown

from backend.core.config import Config
from backend.core.llm_client import LLMClient
from backend.core.doc_converter import doc_to_docx, cleanup_temp_docx
from backend.prompts.builtin import (
    IMAGE_EXTENSIONS,
    AUDIO_EXTENSIONS,
    DOCUMENT_EXTENSIONS,
    PROMPTS,
)


class ConversionResult:
    """Result of a file conversion."""

    def __init__(
        self,
        source_path: Path,
        markdown: str = "",
        title: str = "",
        success: bool = True,
        error: str = "",
        used_llm: bool = False,
        elapsed_ms: int = 0,
        logs: list[str] | None = None,
    ):
        self.source_path = source_path
        self.markdown = markdown
        self.title = title or source_path.stem
        self.success = success
        self.error = error
        self.used_llm = used_llm
        self.elapsed_ms = elapsed_ms
        self.logs = logs or []

    @property
    def char_count(self) -> int:
        return len(self.markdown)

    @property
    def word_count(self) -> int:
        return len(self.markdown.split())


class Converter:
    """Handles file conversion using markitdown and optional LLM analysis."""

    def __init__(self, config: Config, log: Callable[[str], None] | None = None):
        self.config = config
        self._md = MarkItDown()
        self._llm: LLMClient | None = None
        self._log = log or (lambda msg: None)
        self._logs: list[str] = []

    def _emit(self, msg: str):
        """Log to both callback and internal list."""
        self._log(msg)
        self._logs.append(msg)

    @property
    def llm(self) -> LLMClient:
        if self._llm is None:
            self._llm = LLMClient(self.config)
        return self._llm

    # =========================================================================
    # 基础转换（无 LLM）— 同步方法，在 thread pool 中执行
    # =========================================================================
    def convert(self, file_path: str | Path) -> ConversionResult:
        self._logs.clear()
        file_path = Path(file_path)
        if not file_path.exists():
            return ConversionResult(file_path, success=False, error="文件不存在", logs=list(self._logs))

        self._emit(f"正在转换: {file_path.name}")
        start = time.time()

        # Handle legacy .doc files
        tmp_docx = None
        actual_path = file_path
        if file_path.suffix.lower() == ".doc":
            self._emit("检测到 .doc 格式，正在转换为 .docx...")
            tmp_docx = doc_to_docx(file_path)
            if tmp_docx:
                actual_path = tmp_docx
                self._emit(".doc 转换成功，继续处理...")
            else:
                elapsed = int((time.time() - start) * 1000)
                msg = "无法转换 .doc 文件，请确保已安装 Microsoft Word，或将文件另存为 .docx 格式"
                self._emit(f"转换失败: {file_path.name} — {msg}")
                return ConversionResult(
                    file_path, success=False, error=msg, elapsed_ms=elapsed, logs=list(self._logs)
                )

        try:
            result = self._md.convert(str(actual_path))
            markdown = result.text_content if hasattr(result, "text_content") else str(result)
        except Exception as e:
            elapsed = int((time.time() - start) * 1000)
            self._emit(f"转换失败: {file_path.name} — {e}")
            return ConversionResult(
                file_path, success=False, error=str(e), elapsed_ms=elapsed, logs=list(self._logs)
            )
        finally:
            cleanup_temp_docx(tmp_docx)

        elapsed = int((time.time() - start) * 1000)
        self._emit(f"转换完成: {file_path.name} ({len(markdown)} 字符, {elapsed}ms)")
        return ConversionResult(
            file_path,
            markdown=markdown,
            title=file_path.stem,
            elapsed_ms=elapsed,
            logs=list(self._logs),
        )

    # =========================================================================
    # LLM 增强转换（主入口）— 纯 async，由 FastAPI 事件循环驱动
    # =========================================================================
    async def convert_with_llm(self, file_path: str | Path) -> ConversionResult:
        """LLM-enhanced conversion. Pure async — no more prompt_key cruft."""
        self._logs.clear()
        file_path = Path(file_path)
        start = time.time()

        # Handle legacy .doc files
        tmp_docx = None
        actual_path = file_path
        if file_path.suffix.lower() == ".doc":
            self._emit("检测到 .doc 格式，正在转换为 .docx...")
            tmp_docx = doc_to_docx(file_path)
            if tmp_docx:
                actual_path = tmp_docx
                self._emit(".doc 转换成功，继续处理...")

        # Step 1: Base conversion (run sync via thread pool)
        loop = asyncio.get_event_loop()
        base_result = await loop.run_in_executor(None, self.convert, actual_path)

        if not base_result.success:
            cleanup_temp_docx(tmp_docx)
            base_result.source_path = file_path
            return base_result

        base_result.source_path = file_path
        ext = file_path.suffix.lower()
        actual_ext = actual_path.suffix.lower()
        enable_image = self.config.get("conversion.enable_llm_image", False)
        enable_audio = self.config.get("conversion.enable_llm_audio", False)
        enable_summary = self.config.get("conversion.enable_summary", False)
        enable_form_cleaning = self.config.get("conversion.enable_form_cleaning", False)

        self._emit(f"LLM 功能状态 — 图片:{enable_image} 音频:{enable_audio} 总结:{enable_summary} 表单清洗:{enable_form_cleaning}")

        markdown = base_result.markdown

        try:
            # Step 2: Image descriptions
            if enable_image:
                if ext in IMAGE_EXTENSIONS or actual_ext in IMAGE_EXTENSIONS:
                    self._emit("正在分析图片...")
                    desc = await self.llm.analyze_image(
                        actual_path, PROMPTS["image"]["prompt"]
                    )
                    markdown = f"### 图片描述\n\n{desc}\n\n---\n\n{markdown}"
                    base_result.used_llm = True
                    self._emit("图片分析完成")

                elif ext in DOCUMENT_EXTENSIONS or actual_ext in DOCUMENT_EXTENSIONS:
                    markdown = await self._inline_image_descriptions(actual_path, markdown)
                    base_result.used_llm = True

            # Step 3: Audio transcription
            if enable_audio and (ext in AUDIO_EXTENSIONS or actual_ext in AUDIO_EXTENSIONS):
                self._emit("正在转录音频...")
                llm_result = await self.llm.chat(
                    f"{PROMPTS['audio']['prompt']}\n\n---\n\n{markdown[:8000]}"
                )
                markdown = llm_result
                base_result.used_llm = True
                self._emit("音频转录完成")

            # Step 4: Form cleaning
            if enable_form_cleaning and (ext in DOCUMENT_EXTENSIONS or actual_ext in DOCUMENT_EXTENSIONS):
                before_len = len(markdown)
                self._emit(f"正在清洗表单... (文档{ext}, 清洗前{before_len}字符)")
                markdown = await self._clean_form(markdown)
                after_len = len(markdown)
                base_result.used_llm = True
                self._emit(f"表单清洗完成 (清洗前{before_len}→清洗后{after_len}字符, 变化{after_len - before_len:+d})")

            # Step 5: Content summary
            if enable_summary and (ext in DOCUMENT_EXTENSIONS or actual_ext in DOCUMENT_EXTENSIONS):
                self._emit("正在生成内容分析...")
                markdown = await self._add_summary(markdown)
                base_result.used_llm = True
                self._emit("内容分析完成")

        except Exception as e:
            self._emit(f"LLM 处理异常: {e}")
            markdown += f"\n\n---\n\n> LLM 处理异常: {e}"
        finally:
            cleanup_temp_docx(tmp_docx)

        base_result.markdown = markdown
        base_result.elapsed_ms = int((time.time() - start) * 1000)
        return base_result

    # =========================================================================
    # 图片内联描述
    # =========================================================================
    @staticmethod
    def _get_heading_level_at(markdown: str, position: int) -> int:
        before_text = markdown[:position]
        headings = list(re.finditer(r'^(#{1,6})\s+.+$', before_text, re.MULTILINE))
        if not headings:
            return 2
        last = headings[-1]
        return len(last.group(1))

    @staticmethod
    def _heading_prefix(level: int) -> str:
        return "#" * level + " "

    async def _inline_image_descriptions(self, file_path: Path, markdown: str) -> str:
        images = self._extract_images(file_path)
        if not images:
            return markdown

        total = len(images)
        self._emit(f"检测到 {total} 张内嵌图片")

        img_pattern = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')
        matches = list(img_pattern.finditer(markdown))

        if matches:
            for i, match in enumerate(matches):
                if i >= len(images):
                    break
                img_name, img_data = images[i]
                self._emit(f"正在描述图片 ({i+1}/{total})...")
                try:
                    desc = await self._analyze_image_data(img_data, img_name)
                except Exception as e:
                    desc = f"（图片分析失败: {e}）"
                    self._emit(f"图片 {i+1} 分析失败: {e}")

                parent_level = self._get_heading_level_at(markdown, match.start())
                img_heading_level = min(max(parent_level + 1, 3), 6)
                sub_heading_level = min(img_heading_level + 1, 6)
                img_prefix = self._heading_prefix(img_heading_level)
                adjusted_desc = self._adjust_heading_levels(desc, from_level=4, to_level=sub_heading_level)

                alt_text = match.group(1) or "图片"
                replacement = (
                    f"{img_prefix}{alt_text}\n\n"
                    f"> 原文稿此处有一张图片，以下是对该图的描述。\n\n"
                    f"{adjusted_desc}"
                )
                markdown = markdown.replace(match.group(0), replacement, 1)
            self._emit(f"图片描述完成 ({total} 张)")
        else:
            desc_parts = []
            for i, (img_name, img_data) in enumerate(images):
                self._emit(f"正在描述图片 ({i+1}/{total})...")
                try:
                    desc = await self._analyze_image_data(img_data, img_name)
                    adjusted_desc = self._adjust_heading_levels(desc, from_level=4, to_level=4)
                    desc_parts.append(f"#### 图片 {i + 1}（{img_name}）\n\n{adjusted_desc}")
                except Exception as e:
                    desc_parts.append(f"#### 图片 {i + 1}（{img_name}）\n\n（分析失败: {e}）")

            if desc_parts:
                header = f"### 文档内嵌图片（共 {len(images)} 张）\n\n" + "\n\n".join(desc_parts) + "\n\n---\n\n"
                markdown = header + markdown
            self._emit(f"图片描述完成 ({total} 张)")

        return markdown

    @staticmethod
    def _adjust_heading_levels(text: str, from_level: int, to_level: int) -> str:
        if from_level == to_level:
            return text
        diff = to_level - from_level
        def replace_heading(m):
            hashes = m.group(1)
            new_level = len(hashes) + diff
            new_level = max(1, min(new_level, 6))
            return "#" * new_level + " "
        lines = text.split('\n')
        result = []
        for line in lines:
            m = re.match(r'^(#{' + str(from_level) + r',6})\s+', line)
            if m:
                line = replace_heading(m) + line[m.end():]
            result.append(line)
        return '\n'.join(result)

    async def _analyze_image_data(self, img_data: bytes, img_name: str) -> str:
        suffix = Path(img_name).suffix.lower() or ".png"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(img_data)
            tmp_path = tmp.name
        try:
            result = await self.llm.analyze_image(tmp_path, PROMPTS["image"]["prompt"])
            return result.strip()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    # =========================================================================
    # 复杂表单清洗
    # =========================================================================
    CHUNK_THRESHOLD = 12000

    async def _clean_form(self, markdown: str) -> str:
        if not markdown.strip():
            return markdown

        original_len = len(markdown)
        prompt = PROMPTS["form_cleaning"]["prompt"]

        try:
            if original_len <= self.CHUNK_THRESHOLD:
                cleaned = await self._clean_form_chunk(prompt, markdown)
            else:
                chunks = self._smart_split(markdown, self.CHUNK_THRESHOLD)
                total = len(chunks)
                self._emit(f"表单清洗: 文档过长，分为 {total} 块处理")
                cleaned_parts = []
                for i, chunk in enumerate(chunks):
                    if not chunk.strip():
                        continue
                    self._emit(f"正在清洗表单 ({i+1}/{total})...")
                    part = await self._clean_form_chunk(prompt, chunk)
                    cleaned_parts.append(part)
                cleaned = "\n\n".join(cleaned_parts)

            cleaned_len = len(cleaned)
            ratio = cleaned_len / original_len if original_len > 0 else 1.0

            if ratio < 0.5:
                self._emit(f"表单清洗警告: 保留率仅 {ratio:.0%}，已回退")
                return (
                    markdown
                    + "\n\n---\n\n"
                    + "> **表单清洗警告**：LLM 清洗后内容显著减少"
                    f"（原始 {original_len} 字符 → 清洗后 {cleaned_len} 字符，"
                    f"保留率 {ratio:.0%}），已自动回退到原始转换结果。\n"
                    + "> 请检查文档内容后重试。"
                )

            return cleaned

        except Exception as e:
            self._emit(f"表单清洗失败: {e}")
            return markdown + f"\n\n> 表单清洗失败，返回原始内容: {e}"

    async def _clean_form_chunk(self, prompt: str, chunk: str) -> str:
        return await self.llm.chat(
            f"{prompt}\n\n---\n\n以下是需要清洗的原始 Markdown：\n\n{chunk}"
        )

    @staticmethod
    def _find_best_split_level(markdown: str) -> int:
        for level in range(1, 5):
            pattern = re.compile(rf'^#{{{level}}}\s+', re.MULTILINE)
            count = len(pattern.findall(markdown))
            if count >= 2:
                return level
        return 2

    @staticmethod
    def _is_inside_table(lines: list[str], line_idx: int) -> bool:
        for i in range(line_idx, -1, -1):
            stripped = lines[i].strip()
            if not stripped:
                return False
            if stripped.startswith('|'):
                return True
            return False
        return False

    def _smart_split(self, markdown: str, max_chunk_size: int = 12000) -> list[str]:
        if len(markdown) <= max_chunk_size:
            return [markdown]

        split_level = self._find_best_split_level(markdown)
        lines = markdown.split('\n')
        split_points = [0]

        for i, line in enumerate(lines):
            if i == 0:
                continue
            heading_match = re.match(r'^(#{1,6})\s+', line)
            if heading_match:
                level = len(heading_match.group(1))
                if level <= split_level:
                    if not self._is_inside_table(lines, i):
                        split_points.append(i)

        chunks = []
        for idx, start_line in enumerate(split_points):
            end_line = split_points[idx + 1] if idx + 1 < len(split_points) else len(lines)
            chunk = '\n'.join(lines[start_line:end_line])
            chunks.append(chunk)

        merged = []
        for chunk in chunks:
            if merged and len(merged[-1]) + len(chunk) < max_chunk_size:
                merged[-1] += '\n\n' + chunk
            else:
                merged.append(chunk)

        result = []
        for chunk in merged:
            if len(chunk) > max_chunk_size * 1.5:
                sub_chunks = self._smart_split(chunk, max_chunk_size)
                result.extend(sub_chunks)
            else:
                result.append(chunk)

        return result if result else [markdown]

    # =========================================================================
    # 全文内容总结/分析
    # =========================================================================
    async def _add_summary(self, markdown: str) -> str:
        if not markdown.strip():
            return markdown
        prompt = PROMPTS["document_summary"]["prompt"]

        try:
            if len(markdown) <= self.CHUNK_THRESHOLD:
                summary = await self._summary_chunk(prompt, markdown)
            else:
                chunks = self._smart_split(markdown, self.CHUNK_THRESHOLD)
                total = len(chunks)
                self._emit(f"内容分析: 文档过长，分为 {total} 块处理")
                summaries = []
                for i, chunk in enumerate(chunks):
                    if not chunk.strip():
                        continue
                    self._emit(f"正在分析内容 ({i+1}/{total})...")
                    part = await self._summary_chunk(prompt, chunk)
                    summaries.append(part)
                summary = "\n\n".join(summaries)

            return markdown + f"\n\n---\n\n## LLM 内容分析\n\n{summary}"
        except Exception as e:
            self._emit(f"内容分析失败: {e}")
            return markdown + f"\n\n---\n\n> 内容分析失败: {e}"

    async def _summary_chunk(self, prompt: str, chunk: str) -> str:
        return await self.llm.chat(
            f"{prompt}\n\n---\n\n以下是需要分析的文档内容：\n\n{chunk}"
        )

    # =========================================================================
    # 文档内嵌图片提取
    # =========================================================================
    def _extract_images(self, file_path: Path) -> list[tuple[str, bytes]]:
        ext = file_path.suffix.lower()
        if ext in (".docx", ".doc"):
            return self._extract_images_from_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            return self._extract_images_from_pptx(file_path)
        elif ext == ".pdf":
            return self._extract_images_from_pdf(file_path)
        return []

    def _extract_images_from_docx(self, file_path: Path) -> list[tuple[str, bytes]]:
        images = []
        try:
            from docx import Document
            doc = Document(str(file_path))
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    img_data = rel.target_part.blob
                    img_name = rel.target_ref
                    images.append((img_name, img_data))
        except Exception:
            pass
        return images

    def _extract_images_from_pptx(self, file_path: Path) -> list[tuple[str, bytes]]:
        images = []
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        img_data = shape.image.blob
                        img_name = f"slide_{slide.slide_id}_{shape.name}"
                        images.append((img_name, img_data))
        except Exception:
            pass
        return images

    def _extract_images_from_pdf(self, file_path: Path) -> list[tuple[str, bytes]]:
        images = []
        try:
            import pdfplumber
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    for j, img_info in enumerate(page.images):
                        try:
                            x0 = img_info["x0"]
                            y0 = img_info["top"]
                            x1 = img_info["x1"]
                            y1 = img_info["bottom"]
                            cropped = page.crop((x0, y0, x1, y1))
                            pil_img = cropped.to_image(resolution=150).original
                            import io
                            buf = io.BytesIO()
                            pil_img.save(buf, format="PNG")
                            images.append((f"page_{i+1}_img_{j+1}.png", buf.getvalue()))
                        except Exception:
                            pass
        except Exception:
            pass
        return images

